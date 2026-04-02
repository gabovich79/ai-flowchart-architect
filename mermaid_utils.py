"""
AI Flowchart Architect - Mermaid Utilities
Validation, cleaning, styling, and export helpers.
"""

import re
import base64
import io
import urllib.parse
import requests
from config import MERMAID_INK_BASE


def clean_response(text: str) -> str:
    """
    Strip markdown backticks and extraneous text from LLM output.
    Returns only the Mermaid code block.
    """
    text = text.strip()

    # Remove ```mermaid ... ``` wrapper
    pattern = r"```(?:mermaid)?\s*\n?(.*?)```"
    match = re.search(pattern, text, re.DOTALL)
    if match:
        text = match.group(1).strip()

    # Remove any leading text before 'graph' or 'flowchart'
    graph_match = re.search(r"((?:graph|flowchart)\s+(?:TD|TB|BT|LR|RL))", text)
    if graph_match:
        text = text[graph_match.start() :]

    return text.strip()


def validate_mermaid(code: str) -> bool:
    """
    Basic validation that the code looks like valid Mermaid syntax.
    Checks for required structure elements.
    """
    if not code or not code.strip():
        return False

    code_lower = code.strip().lower()

    # Must start with graph or flowchart directive
    if not re.match(r"^(graph|flowchart)\s+(td|tb|bt|lr|rl)", code_lower):
        return False

    # Must have at least one arrow connection
    if "-->" not in code and "---" not in code and "-.->":
        return False

    # Must have at least one node definition
    if "[" not in code and "(" not in code and "{" not in code:
        return False

    return True


def inject_styles(code: str) -> str:
    """
    Auto-inject color styles for different node types if not already present.
    - Decision nodes (diamond {}) -> Orange
    - Start nodes (([...])) -> Green
    - End/terminal nodes (([...])) -> Red (if last)
    - Process nodes ([...]) -> Blue
    """
    # If styles are already present, don't double-inject
    if "style " in code and "fill:" in code:
        return code

    lines = code.strip().split("\n")
    nodes = _extract_nodes(code)
    style_lines = []

    for node_id, node_type in nodes.items():
        if node_type == "decision":
            style_lines.append(
                f"    style {node_id} fill:#FF8C00,stroke:#FF6600,color:#fff"
            )
        elif node_type == "start":
            style_lines.append(
                f"    style {node_id} fill:#2ECC71,stroke:#27AE60,color:#fff"
            )
        elif node_type == "end":
            style_lines.append(
                f"    style {node_id} fill:#E74C3C,stroke:#C0392B,color:#fff"
            )
        elif node_type == "process":
            style_lines.append(
                f"    style {node_id} fill:#3498DB,stroke:#2980B9,color:#fff"
            )

    if style_lines:
        return "\n".join(lines + style_lines)
    return code


def _extract_nodes(code: str) -> dict:
    """
    Parse Mermaid code and classify nodes by type.
    Returns dict of {node_id: node_type}.
    """
    nodes = {}
    # Find all node definitions
    # Decision: ID{...} or ID{...}
    for match in re.finditer(r"(\w+)\s*\{([^}]+)\}", code):
        nodes[match.group(1)] = "decision"

    # Round/stadium: ID([...]) - typically start/end
    for match in re.finditer(r"(\w+)\s*\(\[([^\]]+)\]\)", code):
        node_id = match.group(1)
        label = match.group(2).lower()
        if any(w in label for w in ["start", "begin", "התחלה", "התחל"]):
            nodes[node_id] = "start"
        elif any(w in label for w in ["end", "finish", "סוף", "סיום"]):
            nodes[node_id] = "end"
        elif node_id not in nodes:
            nodes[node_id] = "end"  # Default round nodes to terminal

    # Square brackets: ID[...] - process nodes
    for match in re.finditer(r"(\w+)\s*\[([^\]]+)\]", code):
        node_id = match.group(1)
        # Don't override if already classified
        if node_id not in nodes:
            # Check it's not part of ([...])
            idx = match.start()
            if idx > 0 and code[idx - 1] == "(":
                continue
            nodes[node_id] = "process"

    return nodes


def get_mermaid_ink_url(code: str, fmt: str = "svg") -> str:
    """
    Generate a mermaid.ink URL for the given code.

    Args:
        code: Mermaid diagram code
        fmt: 'svg' or 'img' (png)

    Returns:
        URL string for the rendered diagram
    """
    encoded = base64.urlsafe_b64encode(code.encode("utf-8")).decode("utf-8")
    if fmt == "svg":
        return f"{MERMAID_INK_BASE}/svg/{encoded}"
    else:
        return f"{MERMAID_INK_BASE}/img/{encoded}"


def export_diagram(code: str, fmt: str = "png") -> bytes | None:
    """
    Export diagram as image bytes using mermaid.ink API.

    Args:
        code: Mermaid diagram code
        fmt: 'png' or 'svg'

    Returns:
        Image bytes or None on failure
    """
    try:
        if fmt == "svg":
            url = get_mermaid_ink_url(code, "svg")
        else:
            url = get_mermaid_ink_url(code, "img")

        response = requests.get(url, timeout=15)
        if response.status_code == 200:
            return response.content
        return None
    except Exception:
        return None


def export_to_pptx(code: str) -> bytes | None:
    """
    Export diagram as a PowerPoint (.pptx) file.
    Creates a 16:9 presentation with the flowchart image centered on a slide.

    Args:
        code: Mermaid diagram code

    Returns:
        PPTX file bytes or None on failure
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        # Get PNG image of the diagram
        png_bytes = export_diagram(code, "png")
        if not png_bytes:
            return None

        # Create 16:9 presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── Slide 1: Title slide ──
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide1 = prs.slides.add_slide(blank_layout)

        # Dark background
        bg = slide1.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0E, 0x11, 0x17)

        # Title text
        from pptx.util import Inches, Pt
        txBox = slide1.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "AI Flowchart Architect"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)

        # Subtitle
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = "Auto-generated flowchart"
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        # ── Slide 2: Diagram ──
        slide2 = prs.slides.add_slide(blank_layout)

        # White background for diagram readability
        bg2 = slide2.background
        fill2 = bg2.fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Insert PNG image centered
        img_stream = io.BytesIO(png_bytes)
        # Calculate dimensions to fit within slide with margins
        max_w = Inches(11)
        max_h = Inches(6.5)
        pic = slide2.shapes.add_picture(
            img_stream, Inches(1.167), Inches(0.5), max_w, max_h
        )

        # Add Mermaid source code as speaker notes
        notes_slide = slide2.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"Mermaid Source Code:\n\n{code}"

        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()

    except Exception:
        return None
